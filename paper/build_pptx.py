"""Build the 12-slide CB-SAFE presentation (teacher-approved structure):
Intro -> Background -> one closely-related paper -> proposal -> evaluation KPIs
-> results per KPI -> conclusion. Diagram-first, minimal math.
Run: python build_pptx.py -> cbsafe_slides.pptx"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
FIGS = os.path.join(HERE, "..", "results", "figs")

BLUE = RGBColor(0x2A, 0x78, 0xD6)
AQUA = RGBColor(0x1B, 0xAF, 0x7A)
RED = RGBColor(0xE3, 0x49, 0x48)
INK = RGBColor(0x0B, 0x0B, 0x0B)
BODY = RGBColor(0x26, 0x25, 0x1F)
MUTED = RGBColor(0x89, 0x87, 0x81)
GREY = RGBColor(0x52, 0x51, 0x4E)
BG = RGBColor(0xFC, 0xFC, 0xFB)
PALE = RGBColor(0xF2, 0xF5, 0xFA)
PALE_G = RGBColor(0xEE, 0xF7, 0xF2)
PALE_R = RGBColor(0xFB, 0xEF, 0xEF)
FONT = "Segoe UI"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
blank = prs.slide_layouts[6]


def new_slide():
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def add_text(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
             spacing=1.0):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_bullets(slide, x, y, w, h, bullets, size=16, spacing=1.15, gap=12):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (lead, rest) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = spacing
        p.space_after = Pt(gap)
        b = p.add_run(); b.text = "▪  "
        b.font.name = FONT; b.font.size = Pt(size); b.font.color.rgb = BLUE
        if lead:
            r = p.add_run(); r.text = lead + " "
            r.font.name = FONT; r.font.size = Pt(size); r.font.bold = True
            r.font.color.rgb = INK
        r2 = p.add_run(); r2.text = rest
        r2.font.name = FONT; r2.font.size = Pt(size); r2.font.color.rgb = BODY
    return box


def box(slide, x, y, w, h, text, fill=PALE, edge=BLUE, size=13, bold=False,
        color=INK, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = edge; sp.line.width = Pt(1.25)
    tf = sp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Pt(6)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color
    return sp


def arrow(slide, x, y, w, label=None, h=Inches(0.32), color=GREY, down=False):
    shp = MSO_SHAPE.DOWN_ARROW if down else MSO_SHAPE.RIGHT_ARROW
    sp = slide.shapes.add_shape(shp, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    if label:
        add_text(slide, x - Inches(0.45), y - Inches(0.34), w + Inches(0.9),
                 Inches(0.3), label, 10.5, GREY, align=PP_ALIGN.CENTER)
    return sp


def chrome(slide, kicker, title, n, title_size=30):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.09), SLIDE_H)
    band.fill.solid(); band.fill.fore_color.rgb = BLUE; band.line.fill.background()
    add_text(slide, Inches(0.75), Inches(0.42), Inches(11.8), Inches(0.35),
             kicker.upper(), 12, BLUE, bold=True)
    add_text(slide, Inches(0.75), Inches(0.75), Inches(11.9), Inches(0.95),
             title, title_size, INK, bold=True)
    add_text(slide, Inches(0.75), Inches(7.05), Inches(6), Inches(0.3),
             "CB-SAFE · Suva & Chowdhury · AIUB", 10.5, MUTED)
    add_text(slide, Inches(12.1), Inches(7.05), Inches(0.9), Inches(0.3),
             f"{n} / 12", 10.5, MUTED)


def add_fig(slide, name, cap, x=Inches(7.35), y=Inches(2.0), w=Inches(5.55),
            capy=Inches(6.35)):
    slide.shapes.add_picture(os.path.join(FIGS, name), x, y, width=w)
    add_text(slide, x, capy, w, Inches(0.55), cap, 11, GREY, align=PP_ALIGN.CENTER)


# ============ 1: TITLE ============
s = new_slide()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.09), SLIDE_H)
band.fill.solid(); band.fill.fore_color.rgb = BLUE; band.line.fill.background()
add_text(s, Inches(0.75), Inches(2.1), Inches(11.8), Inches(0.4),
         "RESEARCH PRESENTATION · POST-QUANTUM CRYPTOGRAPHY × FEDERATED LEARNING",
         13, BLUE, bold=True)
add_text(s, Inches(0.75), Inches(2.6), Inches(11.6), Inches(2.2),
         "CB-SAFE: Code-Based Post-Quantum Secure Aggregation for "
         "Byzantine-Robust Federated Learning", 40, INK, bold=True, spacing=1.1)
add_text(s, Inches(0.75), Inches(4.9), Inches(11.6), Inches(0.5),
         "Md Wahiduzzaman Suva (26-94088-2)   ·   Esm-e Moula Chowdhury Abha (26-94089-2)",
         17, BODY)
add_text(s, Inches(0.75), Inches(5.45), Inches(11.6), Inches(0.4),
         "Department of Computer Science, American International University–Bangladesh",
         14, GREY)
add_text(s, Inches(12.1), Inches(7.05), Inches(0.9), Inches(0.3), "1 / 12", 10.5, MUTED)

# ============ 2: INTRODUCTION ============
s = new_slide(); chrome(s, "Introduction", "What this research is about", 2)
add_bullets(s, Inches(0.75), Inches(1.9), Inches(11.9), Inches(2.4), [
    ("The setting:", "federated learning (FL) trains one shared model across many devices; data never leaves the client, only model updates travel."),
    ("What we are trying to do:", "keep those updates private against future quantum computers AND keep training correct when some clients are malicious."),
    ("What we propose:", "CB-SAFE, a secure-aggregation framework whose post-quantum protection comes from a code-based scheme (HQC) instead of the usual lattices, extended with a poisoning defense."),
])
# mini-diagram: FL round
y0 = Inches(4.7)
box(s, Inches(0.9), y0, Inches(2.3), Inches(1.0), "Clients train locally\n(private data)")
arrow(s, Inches(3.35), y0 + Inches(0.34), Inches(0.85), "encrypted updates")
box(s, Inches(4.35), y0, Inches(2.5), Inches(1.0), "Server aggregates\n(sees no individual update)")
arrow(s, Inches(7.0), y0 + Inches(0.34), Inches(0.85), "new global model")
box(s, Inches(8.0), y0, Inches(2.3), Inches(1.0), "Better shared model,\nprivacy intact", fill=PALE_G, edge=AQUA)
box(s, Inches(10.8), y0, Inches(1.9), Inches(1.0), "Threats:\nquantum + poisoning", fill=PALE_R, edge=RED)

# ============ 3: BACKGROUND ============
s = new_slide(); chrome(s, "Background", "The problem: three threats meet in one protocol", 3)
add_bullets(s, Inches(0.75), Inches(1.85), Inches(11.9), Inches(2.6), [
    ("Updates leak.", "Gradient-inversion attacks reconstruct private training images from a single client's update, so updates must be hidden."),
    ("Secure aggregation hides them", "with pairwise masks that cancel in the sum, but its key exchange breaks once a quantum computer exists (“harvest now, decrypt later”)."),
    ("All post-quantum fixes use one math family (lattices)", "and hiding updates also blinds the server to poisoned updates from malicious clients."),
])
# threat diagram: three boxes converge on FL aggregation
y0 = Inches(4.55)
box(s, Inches(1.2), y0, Inches(3.0), Inches(0.95), "Privacy threat:\nupdates reveal training data", fill=PALE_R, edge=RED)
box(s, Inches(5.1), y0, Inches(3.0), Inches(0.95), "Quantum threat:\nkey exchange breaks (Shor)", fill=PALE_R, edge=RED)
box(s, Inches(9.0), y0, Inches(3.0), Inches(0.95), "Poisoning threat:\nmalicious clients flip the model", fill=PALE_R, edge=RED)
for cx in (2.7, 6.6, 10.5):
    arrow(s, Inches(cx - 0.16), y0 + Inches(1.05), Inches(0.32), down=True,
          h=Inches(0.4))
box(s, Inches(3.4), y0 + Inches(1.55), Inches(6.4), Inches(0.85),
    "FL secure aggregation must survive all three at once", bold=True, size=14)

# ============ 4: RELATED WORK ============
s = new_slide(); chrome(s, "Closely related work",
                        "How the closest peer-reviewed paper solves it", 4)
add_text(s, Inches(0.75), Inches(1.8), Inches(11.9), Inches(0.6),
         "PQSF: “Post-quantum secure privacy-preserving federated learning”, "
         "Zhang et al., Scientific Reports 14, 2024 (peer-reviewed)", 15.5, GREY, bold=True)
add_bullets(s, Inches(0.75), Inches(2.45), Inches(11.9), Inches(2.2), [
    ("Their solution:", "the Bonawitz-style masking protocol, made quantum-safe with lattice cryptography: a lattice multi-stage secret-sharing scheme reconstructs the masks, so the server still learns only the sum of updates."),
    ("What works:", "peer-reviewed evidence that post-quantum secure aggregation is practical in FL; our confidentiality layer follows the same masking pattern."),
])
# gap boxes
y0 = Inches(4.65)
box(s, Inches(0.9), y0, Inches(3.7), Inches(1.15),
    "Gap 1 · Lattice-only\nOne cryptanalytic advance\nbreaks everything", fill=PALE_R, edge=RED, size=12.5)
box(s, Inches(4.85), y0, Inches(3.7), Inches(1.15),
    "Gap 2 · No robustness\nA poisoned update passes\nstraight through the sum", fill=PALE_R, edge=RED, size=12.5)
box(s, Inches(8.8), y0, Inches(3.7), Inches(1.15),
    "Gap 3 · Hard-wired primitive\nCannot swap the KEM\nwhen standards evolve", fill=PALE_R, edge=RED, size=12.5)
add_text(s, Inches(0.9), y0 + Inches(1.3), Inches(11.6), Inches(0.6),
         "CB-SAFE targets exactly these gaps. (2026 preprints confirm the gap is live: "
         "lattice+reputation and quantum-hardware approaches, but none is code-based.)",
         13, BLUE, bold=True)

# ============ 5: PROPOSAL OVERVIEW ============
s = new_slide(); chrome(s, "Our proposal", "CB-SAFE: one training round, end to end", 5)
add_text(s, Inches(0.75), Inches(1.75), Inches(11.9), Inches(0.5),
         "Code-based (HQC) key setup happens once; every round then runs entirely on cheap masks.",
         15, BODY)
y0 = Inches(2.6)
# clusters column
box(s, Inches(0.85), y0, Inches(2.15), Inches(0.8), "Cluster 1\n(c clients, masked)")
box(s, Inches(0.85), y0 + Inches(1.0), Inches(2.15), Inches(0.8), "Cluster 2\n(c clients, masked)")
add_text(s, Inches(0.85), y0 + Inches(1.9), Inches(2.15), Inches(0.4), "⋮   (k clusters)",
         16, GREY, align=PP_ALIGN.CENTER)
arrow(s, Inches(3.15), y0 + Inches(0.95), Inches(0.8), "masks cancel")
box(s, Inches(4.1), y0 + Inches(0.5), Inches(2.6), Inches(1.7),
    "Server sees ONLY\nk cluster sums\nS₁ … Sₖ\n(never an individual)", bold=True, size=13.5)
arrow(s, Inches(6.85), y0 + Inches(0.95), Inches(0.8), "inspect sums")
box(s, Inches(7.8), y0 + Inches(0.5), Inches(2.5), Inches(1.7),
    "Robust aggregation\n+ CB-SAFE+ defense\n(flag & exclude attackers)", fill=PALE_G, edge=AQUA, size=13.5)
arrow(s, Inches(10.45), y0 + Inches(0.95), Inches(0.8))
box(s, Inches(11.3), y0 + Inches(0.65), Inches(1.6), Inches(1.4), "Global\nmodel", bold=True)
# setup ribbon + broadcast
box(s, Inches(0.85), Inches(5.75), Inches(6.0), Inches(0.7),
    "ONE-TIME: HQC (code-based KEM) key setup (pluggable: swap in ML-KEM/Kyber with one config line)",
    size=12)
box(s, Inches(7.15), Inches(5.75), Inches(5.75), Inches(0.7),
    "EVERY ROUND: fresh masks from stored secrets: zero KEM bytes on the wire",
    size=12, fill=PALE_G, edge=AQUA)

# ============ 6: PROPOSAL - PRIVACY LAYER ============
s = new_slide(); chrome(s, "Proposal · confidentiality layer",
                        "Hide every update, pay for crypto only once", 6)
add_bullets(s, Inches(0.75), Inches(1.9), Inches(11.9), Inches(1.9), [
    ("Each client adds secret masks to its update;", "within a cluster the masks cancel, so the sum is exact but individuals stay hidden (anonymity set = cluster size)."),
    ("Guarantee (proved):", "the server, even colluding with some clients, learns nothing beyond each cluster's sum."),
    ("Dropouts handled:", "secret-sharing recovers the masks of clients who disconnect mid-round."),
])
y0 = Inches(4.35)
box(s, Inches(0.9), y0, Inches(3.4), Inches(1.5),
    "SETUP (once)\npublish HQC public keys,\nestablish pairwise secrets", bold=False, size=13)
arrow(s, Inches(4.45), y0 + Inches(0.6), Inches(0.9), "stored secrets")
box(s, Inches(5.5), y0, Inches(3.4), Inches(1.5),
    "EVERY ROUND\nderive fresh masks locally\n(no new crypto traffic)", fill=PALE_G, edge=AQUA, size=13)
arrow(s, Inches(9.05), y0 + Inches(0.6), Inches(0.9), "masked update")
box(s, Inches(10.1), y0, Inches(2.6), Inches(1.5),
    "SERVER\nsums per cluster,\nmasks vanish", size=13)

# ============ 7: PROPOSAL - DEFENSE ============
s = new_slide(); chrome(s, "Proposal · CB-SAFE+ defense",
                        "Catch attackers the server never sees individually", 7)
add_bullets(s, Inches(0.75), Inches(1.9), Inches(11.9), Inches(1.55), [
    ("Key insight:", "a poisoned cluster sum looks normal in size but pushes the model the wrong way; a tiny server-side dataset (200 samples) exposes it: applying that sum raises the loss."),
    ("Repetition identifies individuals:", "clusters are reshuffled every round, so a real attacker is inside a flagged cluster almost every round while honest clients are flagged only occasionally."),
])
y0 = Inches(3.85)
steps = [
    ("1 · Reshuffle", "new random clusters\nevery round", PALE, BLUE),
    ("2 · Probe", "does each cluster sum\nraise root-set loss?", PALE, BLUE),
    ("3 · Flag", "mark suspicious\ncluster sums", PALE_R, RED),
    ("4 · Score", "suspicion += 1 for all\nmembers of flagged clusters", PALE, BLUE),
    ("5 · Exclude", "persistent suspects are\nremoved from training", PALE_G, AQUA),
]
x = 0.8
for i, (t1, t2, fill, edge) in enumerate(steps):
    box(s, Inches(x), y0, Inches(2.05), Inches(1.45), f"{t1}\n{t2}", fill=fill,
        edge=edge, size=11.5)
    if i < 4:
        arrow(s, Inches(x + 2.1), y0 + Inches(0.57), Inches(0.32), h=Inches(0.3))
    x += 2.45
add_text(s, Inches(0.8), Inches(5.65), Inches(11.9), Inches(0.5),
         "Flags are computed from sums the server already holds, so the defense leaks nothing extra (proved).",
         14, GREY)

# ============ 8: EVALUATION METHODOLOGY / KPIs ============
s = new_slide(); chrome(s, "Evaluation methodology",
                        "How we evaluate: five KPI families", 8)
rows = [
    ("KPI", "What it measures", "How"),
    ("Communication overhead", "bytes per client: one-time setup vs every round", "counted from real protocol messages (liboqs)"),
    ("Computation time", "key setup, masking, unmasking wall-clock", "measured per round on real hardware"),
    ("Utility", "test accuracy vs plain FedAvg (must match)", "CIFAR-10, non-IID, 30 clients, 40 rounds"),
    ("Robustness", "accuracy & attack success rate vs attacker fraction f", "sign-flip, label-flip, backdoor at f = 10–30%"),
    ("Identification quality", "attackers caught / honest clients wrongly excluded", "vs ground truth across all runs"),
]
tbl = s.shapes.add_table(len(rows), 3, Inches(0.75), Inches(1.95), Inches(11.85),
                         Inches(3.9)).table
tbl.columns[0].width = Inches(3.0)
tbl.columns[1].width = Inches(4.55)
tbl.columns[2].width = Inches(4.3)
for ri, (a, b, c) in enumerate(rows):
    for ci, val in enumerate((a, b, c)):
        cell = tbl.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE if ri == 0 else (PALE if ri % 2 else BG)
        p = cell.text_frame.paragraphs[0]
        r = p.add_run(); r.text = val
        r.font.name = FONT; r.font.size = Pt(13 if ri else 14)
        r.font.bold = (ri == 0 or ci == 0)
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if ri == 0 else BODY
add_text(s, Inches(0.75), Inches(6.15), Inches(11.9), Inches(0.5),
         "Baselines: plain FedAvg, lattice (ML-KEM) variant, and four standard robust rules "
         "(mean, trimmed mean, median, Multi-Krum), all at matched NIST security levels.",
         13.5, GREY)

# ============ 9: RESULTS - COST KPI ============
s = new_slide(); chrome(s, "Evaluation · KPI 1 & 2 (cost)",
                        "The code-based premium is a one-time constant", 9)
add_bullets(s, Inches(0.75), Inches(2.0), Inches(6.2), Inches(4.4), [
    ("Every round costs the same for all six KEMs:", "2,277 KiB per client; the update itself dominates."),
    ("Setup (once):", "HQC 15.2 KiB vs ML-KEM 3.8 KiB per client, under 0.03% of 40 rounds' traffic."),
    ("Time:", "2.0 s one-time key setup for 30 clients; ~26 ms masking per round vs ~21 s of training."),
    ("Takeaway:", "cryptographic diversity is practically free once keys are reused."),
], size=15.5)
add_fig(s, "fig_amortization.png",
        "Even the heaviest HQC setup is ~2% of ONE round's traffic (log scale).")

# ============ 10: RESULTS - UTILITY KPI ============
s = new_slide(); chrome(s, "Evaluation · KPI 3 (utility)",
                        "Privacy provably costs no accuracy", 10)
add_bullets(s, Inches(0.75), Inches(2.0), Inches(6.2), Inches(4.4), [
    ("59.0% on non-IID CIFAR-10", "(30 clients, 40 rounds), identical for plain, HQC, and ML-KEM runs."),
    ("Why identical:", "secure sums matched plain sums in every round to 2.28×10⁻⁵, exactly the rounding limit; the crypto changes nothing else."),
    ("Dropout recovery verified live", "in rounds 10 and 20, including graceful exclusion of a below-threshold cluster."),
], size=15.5)
add_fig(s, "fig_utility.png",
        "One accuracy curve serves all three configurations.", y=Inches(2.3))

# ============ 11: RESULTS - ROBUSTNESS KPI ============
s = new_slide(); chrome(s, "Evaluation · KPI 4 & 5 (robustness)",
                        "Standard rules collapse; CB-SAFE+ recovers", 11)
add_bullets(s, Inches(0.75), Inches(2.0), Inches(6.2), Inches(4.4), [
    ("Why rules fail:", "averaging inside a cluster disguises the poison (the sum has a normal size, just the wrong direction), so trimmed mean and median fall to 10% (random guessing) at f ≥ 20%."),
    ("CB-SAFE+ result:", "catches 8 of 9 attackers at f = 30% with zero honest clients wrongly excluded, recovering 42.8% where every static rule reads 10–17%."),
    ("Honest limits:", "stealthy backdoors evade all sum-level defenses (93–97% attack success for every rule); Multi-Krum stays stronger at f = 20%."),
], size=15.5)
add_fig(s, "fig_signflip_acc.png",
        "Sign-flip attack: accuracy vs attacker fraction, per aggregation rule.")

# ============ 12: CONCLUSION ============
s = new_slide(); chrome(s, "Conclusion", "Takeaways & what's next", 12)
add_bullets(s, Inches(0.75), Inches(2.0), Inches(11.9), Inches(4.6), [
    ("Diversity delivered:", "the first code-based post-quantum secure aggregation for FL, at a one-time cost under 0.03% of traffic."),
    ("New understanding:", "secure aggregation launders poisoning; defenses must use signals that survive summation, accumulated over rounds."),
    ("New defense:", "CB-SAFE+ identifies individual attackers from group observations only, with zero extra privacy leakage."),
    ("Open problems:", "stealthy backdoors, actively malicious servers, adaptive attackers."),
    ("In progress:", "multi-seed statistics, a second dataset, 100-client scale; implementation released open-source."),
], size=17, gap=14)

out = os.path.join(HERE, "cbsafe_slides.pptx")
prs.save(out)
print("wrote", out, f"({len(prs.slides._sldIdLst)} slides)")

